from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _shape_text(shape) -> str:
    """Return all text from a shape's text frame, or empty string."""
    if shape.has_text_frame:
        return "\n".join(
            para.text.strip()
            for para in shape.text_frame.paragraphs
            if para.text.strip()
        )
    return ""


def _table_chunk(table, slide_number: int) -> dict:
    """Convert a pptx Table object to a structured TABLE chunk."""
    rows = table.rows
    if not rows:
        return {}
    headers = [cell.text.strip() for cell in rows[0].cells]
    lines = ["TABLE: " + " | ".join(headers)]
    for row in list(rows)[1:]:
        for h, cell in zip(headers, row.cells):
            lines.append(f"{h}: {cell.text.strip()}")
        lines.append("---")
    return {"type": "table", "text": "\n".join(lines), "slide_number": slide_number}


def parse_pptx(file_path: str) -> list[dict]:
    """
    Parse a PowerPoint file slide by slide.
    Each slide produces one text chunk (title + body + speaker notes)
    and one chunk per table found on that slide.

    Returns a list of dicts with keys:
        type          — 'text' or 'table'
        text          — raw string content
        slide_number  — 1-based slide number
    """
    prs = Presentation(file_path)
    chunks: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        slide_title = ""

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                pass  # handled after text pass so we have the title
            else:
                t = _shape_text(shape)
                if t:
                    # Use the first text shape (usually the title) as context
                    if not slide_title:
                        slide_title = t.splitlines()[0]
                    text_parts.append(t)

        # Append speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                text_parts.append(f"[Notes] {notes_text}")

        if text_parts:
            chunks.append({
                "type": "text",
                "text": "\n".join(text_parts),
                "slide_number": slide_num,
            })

        # Now emit table chunks with slide context prepended
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                table_chunk = _table_chunk(shape.table, slide_num)
                if table_chunk:
                    if slide_title:
                        table_chunk["text"] = (
                            f"[Context: Slide {slide_num} — {slide_title}]\n\n"
                            + table_chunk["text"]
                        )
                    chunks.append(table_chunk)

    return chunks
