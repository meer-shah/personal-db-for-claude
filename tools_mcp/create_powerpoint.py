"""
create_presentation tool — build a real .pptx from structured content
using python-pptx, upload to OneDrive.
"""

import io

from fastapi import APIRouter, Depends, Request
from pydantic import BaseModel, Field
from slowapi import Limiter
from slowapi.util import get_remote_address

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from tools_mcp.auth import require_bearer
from tools_mcp._onedrive_upload import upload_bytes

router  = APIRouter()
limiter = Limiter(key_func=get_remote_address)

_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class PptSlide(BaseModel):
    title:    str | None = None
    subtitle: str | None = None  # only used on title slide
    bullets:  list[str] = Field(default_factory=list)
    body:     str | None = None  # free-form body text (used if no bullets)
    notes:    str | None = None


class CreatePptRequest(BaseModel):
    filename:    str
    folder_path: str
    title:       str | None = None
    subtitle:    str | None = None
    slides:      list[PptSlide] = Field(default_factory=list)


class CreatePptResponse(BaseModel):
    success:      bool
    onedrive_url: str | None = None
    file_id:      str | None = None
    error:        str | None = None


def _add_title_slide(prs: Presentation, title: str, subtitle: str | None) -> None:
    layout = prs.slide_layouts[0]  # Title Slide
    slide  = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_content_slide(prs: Presentation, sl: PptSlide) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide  = prs.slides.add_slide(layout)

    if slide.shapes.title is not None and sl.title:
        slide.shapes.title.text = sl.title

    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            body_ph = ph
            break

    if body_ph is not None:
        tf = body_ph.text_frame
        tf.word_wrap = True

        if sl.bullets:
            tf.text = sl.bullets[0]
            for bullet in sl.bullets[1:]:
                p = tf.add_paragraph()
                p.text  = bullet
                p.level = 0
        elif sl.body:
            tf.text = sl.body

    if sl.notes:
        slide.notes_slide.notes_text_frame.text = sl.notes


def _build_pptx(req: CreatePptRequest) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if req.title:
        _add_title_slide(prs, req.title, req.subtitle)

    for sl in req.slides:
        _add_content_slide(prs, sl)

    if not prs.slides:
        # Always produce at least one slide
        _add_title_slide(prs, req.title or "Untitled", req.subtitle)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@router.post("/tools/create_presentation", response_model=CreatePptResponse)
@limiter.limit("10/minute")
def create_presentation(
    request: Request,
    req: CreatePptRequest,
    _: None = Depends(require_bearer),
) -> CreatePptResponse:
    try:
        filename = req.filename if req.filename.lower().endswith(".pptx") else f"{req.filename}.pptx"
        raw = _build_pptx(req)
        data = upload_bytes(raw, filename, req.folder_path, _PPTX_MIME)
        return CreatePptResponse(
            success      = True,
            onedrive_url = data.get("webUrl"),
            file_id      = data.get("id"),
        )
    except Exception as e:
        return CreatePptResponse(success=False, error=str(e))
